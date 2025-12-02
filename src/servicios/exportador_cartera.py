"""
Exportador de Cartera de Proyectos - Reportes Profesionales
Sistema de Priorización ENLAZA GEB - Arquitectura C

Genera reportes profesionales para comités de aprobación:
- Word: Documento completo con portada, metodología, ranking y recomendaciones
- Excel: Múltiples hojas con datos, análisis y comparativos
- PowerPoint: Presentación ejecutiva para comité de aprobación

Autor: Sistema ENLAZA
Fecha: Diciembre 2025
"""

from docx import Document
from docx.shared import Inches, Pt, RGBColor, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.enum.style import WD_STYLE_TYPE
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils.dataframe import dataframe_to_rows
from openpyxl.chart import BarChart, Reference, PieChart
from openpyxl.chart.label import DataLabelList
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor as PptxRGBColor
import pandas as pd
from datetime import datetime
from typing import Dict, List, Any, Optional
from io import BytesIO


class ExportadorCartera:
    """
    Exportador profesional de cartera de proyectos.

    Genera documentos Word y Excel con formato profesional
    para presentación ante comités de aprobación.
    """

    # Pesos de criterios Arquitectura C
    PESOS_CRITERIOS = {
        'SROI': 0.40,
        'Stakeholders': 0.25,
        'Probabilidad': 0.20,
        'Riesgos': 0.15
    }

    # Colores corporativos
    COLOR_PRIMARIO = RGBColor(0, 102, 153)  # Azul corporativo
    COLOR_SECUNDARIO = RGBColor(0, 153, 102)  # Verde
    COLOR_ALERTA = RGBColor(204, 0, 0)  # Rojo
    COLOR_ADVERTENCIA = RGBColor(255, 153, 0)  # Naranja

    def __init__(self, reporte: Dict[str, Any], resultados_detallados: List[Any]):
        """
        Inicializa el exportador.

        Args:
            reporte: Diccionario con resumen de la evaluación
            resultados_detallados: Lista de resultados por proyecto (objetos o dicts)
        """
        self.reporte = reporte
        # Normalizar resultados a diccionarios
        self.resultados = self._normalizar_resultados(resultados_detallados)
        self.fecha_generacion = datetime.now()

    def _normalizar_resultados(self, resultados: List[Any]) -> List[Dict[str, Any]]:
        """
        Normaliza resultados a formato diccionario consistente.

        Maneja tanto objetos ResultadoEvaluacion como diccionarios.
        """
        normalizados = []

        for r in resultados:
            # Si es un objeto con método to_dict, usarlo
            if hasattr(r, 'to_dict'):
                dato = r.to_dict()
            elif hasattr(r, '__dict__'):
                # Si es un objeto dataclass u otro objeto
                dato = {
                    'nombre': getattr(r, 'proyecto_nombre', getattr(r, 'nombre', 'Sin nombre')),
                    'proyecto': getattr(r, 'proyecto_nombre', getattr(r, 'nombre', 'Sin nombre')),
                    'score_final': getattr(r, 'score_final', 0),
                    'detalle_criterios': getattr(r, 'detalle_criterios', {}),
                    'alertas': getattr(r, 'observaciones', []),
                    'recomendacion': getattr(r, 'recomendacion', '')
                }
            elif isinstance(r, dict):
                dato = r
            else:
                continue

            # Normalizar nombre
            if 'nombre' not in dato:
                dato['nombre'] = dato.get('proyecto_nombre', dato.get('proyecto', 'Sin nombre'))

            # Normalizar detalle_criterios si es dict de dicts (formato ResultadoEvaluacion)
            if 'detalle_criterios' in dato and isinstance(dato['detalle_criterios'], dict):
                detalle = dato['detalle_criterios']
                # Verificar si es formato {criterio: {score_base, peso, contribucion}}
                if detalle and not isinstance(list(detalle.values())[0] if detalle else {}, dict):
                    pass  # Ya está en formato lista
                else:
                    # Convertir de dict a lista
                    detalle_lista = []
                    for criterio, valores in detalle.items():
                        if isinstance(valores, dict):
                            detalle_lista.append({
                                'criterio': criterio,
                                'score_base': valores.get('score_base', 0),
                                'contribucion_parcial': valores.get('contribucion', valores.get('contribucion_parcial', 0))
                            })
                    dato['detalle_criterios'] = detalle_lista

            # Normalizar alertas
            if 'alertas' not in dato:
                dato['alertas'] = dato.get('observaciones', [])

            normalizados.append(dato)

        return normalizados

    def exportar_word(self) -> BytesIO:
        """
        Genera documento Word profesional con reporte completo.

        Returns:
            BytesIO con el documento Word
        """
        doc = Document()

        # Configurar estilos
        self._configurar_estilos_word(doc)

        # 1. Portada
        self._agregar_portada_word(doc)

        # 2. Resumen Ejecutivo
        self._agregar_resumen_ejecutivo_word(doc)

        # 3. Metodología
        self._agregar_metodologia_word(doc)

        # 4. Ranking de Proyectos
        self._agregar_ranking_word(doc)

        # 5. Análisis Comparativo
        self._agregar_comparativo_word(doc)

        # 6. Detalle por Proyecto
        self._agregar_detalle_proyectos_word(doc)

        # 7. Recomendaciones
        self._agregar_recomendaciones_word(doc)

        # Guardar
        buffer = BytesIO()
        doc.save(buffer)
        buffer.seek(0)
        return buffer

    def exportar_excel(self) -> BytesIO:
        """
        Genera archivo Excel con múltiples hojas de análisis.

        Returns:
            BytesIO con el archivo Excel
        """
        wb = Workbook()

        # 1. Hoja Resumen
        self._crear_hoja_resumen_excel(wb)

        # 2. Hoja Ranking
        self._crear_hoja_ranking_excel(wb)

        # 3. Hoja Comparativo por Criterio
        self._crear_hoja_comparativo_excel(wb)

        # 4. Hoja Detalle
        self._crear_hoja_detalle_excel(wb)

        # 5. Hoja Metodología
        self._crear_hoja_metodologia_excel(wb)

        # Eliminar hoja por defecto si existe
        if 'Sheet' in wb.sheetnames:
            del wb['Sheet']

        # Guardar
        buffer = BytesIO()
        wb.save(buffer)
        buffer.seek(0)
        return buffer

    # ==================== MÉTODOS WORD ====================

    def _configurar_estilos_word(self, doc: Document):
        """Configura estilos del documento Word."""
        # Estilo para títulos
        styles = doc.styles

        # Modificar Heading 1
        style_h1 = styles['Heading 1']
        style_h1.font.color.rgb = self.COLOR_PRIMARIO
        style_h1.font.size = Pt(16)
        style_h1.font.bold = True

        # Modificar Heading 2
        style_h2 = styles['Heading 2']
        style_h2.font.color.rgb = self.COLOR_PRIMARIO
        style_h2.font.size = Pt(14)
        style_h2.font.bold = True

        # Modificar Normal
        style_normal = styles['Normal']
        style_normal.font.size = Pt(11)
        style_normal.font.name = 'Calibri'

    def _agregar_portada_word(self, doc: Document):
        """Agrega portada profesional al documento."""
        # Espaciado superior
        for _ in range(3):
            doc.add_paragraph()

        # Logo/Título empresa
        titulo_empresa = doc.add_paragraph()
        titulo_empresa.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = titulo_empresa.add_run("ENLAZA GEB")
        run.bold = True
        run.font.size = Pt(28)
        run.font.color.rgb = self.COLOR_PRIMARIO

        # Subtítulo
        subtitulo = doc.add_paragraph()
        subtitulo.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = subtitulo.add_run("Sistema de Priorización de Proyectos Sociales")
        run.font.size = Pt(16)
        run.font.color.rgb = RGBColor(100, 100, 100)

        doc.add_paragraph()
        doc.add_paragraph()

        # Título del reporte
        titulo_reporte = doc.add_paragraph()
        titulo_reporte.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = titulo_reporte.add_run("REPORTE DE EVALUACIÓN DE CARTERA")
        run.bold = True
        run.font.size = Pt(24)
        run.font.color.rgb = self.COLOR_PRIMARIO

        doc.add_paragraph()

        # Información de la evaluación
        info = doc.add_paragraph()
        info.alignment = WD_ALIGN_PARAGRAPH.CENTER

        num_proyectos = len(self.resultados)
        fecha = self.fecha_generacion.strftime("%d de %B de %Y")

        run = info.add_run(f"Evaluación de {num_proyectos} proyectos")
        run.font.size = Pt(14)
        info.add_run("\n")
        run = info.add_run(f"Metodología: Arquitectura C")
        run.font.size = Pt(12)
        run.font.italic = True

        # Espaciado inferior
        for _ in range(8):
            doc.add_paragraph()

        # Fecha
        fecha_p = doc.add_paragraph()
        fecha_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = fecha_p.add_run(f"Fecha de generación: {fecha}")
        run.font.size = Pt(11)
        run.font.color.rgb = RGBColor(100, 100, 100)

        # Salto de página
        doc.add_page_break()

    def _agregar_resumen_ejecutivo_word(self, doc: Document):
        """Agrega sección de resumen ejecutivo."""
        doc.add_heading("1. Resumen Ejecutivo", level=1)

        # Métricas principales
        doc.add_heading("1.1 Métricas Principales", level=2)

        # Calcular métricas
        num_proyectos = len(self.resultados)
        if num_proyectos > 0:
            scores = [r.get('score_final', 0) for r in self.resultados]
            score_promedio = sum(scores) / len(scores)
            score_max = max(scores)
            score_min = min(scores)

            # Distribución por prioridad
            alta = sum(1 for r in self.resultados if r.get('score_final', 0) >= 70)
            media = sum(1 for r in self.resultados if 50 <= r.get('score_final', 0) < 70)
            baja = sum(1 for r in self.resultados if r.get('score_final', 0) < 50)
        else:
            score_promedio = score_max = score_min = 0
            alta = media = baja = 0

        # Tabla de métricas
        table = doc.add_table(rows=5, cols=2)
        table.style = 'Table Grid'

        metricas = [
            ("Total de Proyectos Evaluados", str(num_proyectos)),
            ("Score Promedio", f"{score_promedio:.1f}/100"),
            ("Score Máximo", f"{score_max:.1f}/100"),
            ("Score Mínimo", f"{score_min:.1f}/100"),
            ("Proyectos Prioridad Alta (≥70)", f"{alta} ({alta/num_proyectos*100:.0f}%)" if num_proyectos > 0 else "0")
        ]

        for i, (metrica, valor) in enumerate(metricas):
            row = table.rows[i]
            row.cells[0].text = metrica
            row.cells[1].text = valor
            # Formato
            row.cells[0].paragraphs[0].runs[0].bold = True

        doc.add_paragraph()

        # Distribución de prioridades
        doc.add_heading("1.2 Distribución por Prioridad", level=2)

        p = doc.add_paragraph()
        p.add_run(f"• Prioridad ALTA (≥70): ").bold = True
        p.add_run(f"{alta} proyectos")

        p = doc.add_paragraph()
        p.add_run(f"• Prioridad MEDIA (50-69): ").bold = True
        p.add_run(f"{media} proyectos")

        p = doc.add_paragraph()
        p.add_run(f"• Prioridad BAJA (<50): ").bold = True
        p.add_run(f"{baja} proyectos")

        doc.add_paragraph()

        # Top 3 proyectos
        if num_proyectos > 0:
            doc.add_heading("1.3 Top 3 Proyectos Recomendados", level=2)

            # Ordenar por score
            proyectos_ordenados = sorted(
                self.resultados,
                key=lambda x: x.get('score_final', 0),
                reverse=True
            )[:3]

            for i, proyecto in enumerate(proyectos_ordenados, 1):
                nombre = proyecto.get('nombre', proyecto.get('proyecto', 'Sin nombre'))
                score = proyecto.get('score_final', 0)

                p = doc.add_paragraph()
                p.add_run(f"{i}. {nombre}").bold = True
                p.add_run(f" — Score: {score:.1f}/100")

        doc.add_page_break()

    def _agregar_metodologia_word(self, doc: Document):
        """Agrega sección de metodología."""
        doc.add_heading("2. Metodología de Evaluación", level=1)

        doc.add_paragraph(
            "La evaluación utiliza la Arquitectura C del Sistema de Priorización ENLAZA, "
            "diseñada específicamente para proyectos de valor compartido en el sector energético."
        )

        doc.add_heading("2.1 Criterios de Evaluación", level=2)

        # Tabla de criterios
        table = doc.add_table(rows=5, cols=3)
        table.style = 'Table Grid'

        # Encabezados
        headers = table.rows[0].cells
        headers[0].text = "Criterio"
        headers[1].text = "Peso"
        headers[2].text = "Descripción"
        for cell in headers:
            cell.paragraphs[0].runs[0].bold = True

        criterios = [
            ("SROI", "40%", "Retorno Social de la Inversión - Métrica dominante"),
            ("Stakeholders", "25%", "Relacionamiento y pertinencia operacional"),
            ("Probabilidad", "20%", "Probabilidad de éxito del proyecto"),
            ("Riesgos", "15%", "Evaluación de riesgos (scoring inverso)")
        ]

        for i, (criterio, peso, desc) in enumerate(criterios, 1):
            row = table.rows[i].cells
            row[0].text = criterio
            row[1].text = peso
            row[2].text = desc

        doc.add_paragraph()

        doc.add_heading("2.2 Escalas de Priorización", level=2)

        p = doc.add_paragraph()
        p.add_run("• Score ≥ 70: ").bold = True
        p.add_run("Prioridad ALTA - Recomendado para aprobación inmediata")

        p = doc.add_paragraph()
        p.add_run("• Score 50-69: ").bold = True
        p.add_run("Prioridad MEDIA - Considerar con ajustes")

        p = doc.add_paragraph()
        p.add_run("• Score < 50: ").bold = True
        p.add_run("Prioridad BAJA - Requiere rediseño o reconsideración")

        doc.add_paragraph()

        doc.add_heading("2.3 Consideraciones Especiales", level=2)

        consideraciones = [
            "SROI < 1.0: Proyecto rechazado automáticamente (destruye valor)",
            "SROI > 7.0: Requiere verificación metodológica",
            "Riesgos críticos: Alertas especiales independiente del score",
            "Municipios PDET: Bonificación por impacto territorial"
        ]

        for c in consideraciones:
            p = doc.add_paragraph(c, style='List Bullet')

        doc.add_page_break()

    def _agregar_ranking_word(self, doc: Document):
        """Agrega tabla de ranking de proyectos."""
        doc.add_heading("3. Ranking de Proyectos", level=1)

        if not self.resultados:
            doc.add_paragraph("No hay proyectos para mostrar.")
            return

        # Ordenar por score
        proyectos_ordenados = sorted(
            self.resultados,
            key=lambda x: x.get('score_final', 0),
            reverse=True
        )

        # Tabla de ranking
        table = doc.add_table(rows=len(proyectos_ordenados) + 1, cols=4)
        table.style = 'Table Grid'

        # Encabezados
        headers = table.rows[0].cells
        headers[0].text = "Posición"
        headers[1].text = "Proyecto"
        headers[2].text = "Score"
        headers[3].text = "Prioridad"
        for cell in headers:
            cell.paragraphs[0].runs[0].bold = True

        # Datos
        for i, proyecto in enumerate(proyectos_ordenados, 1):
            row = table.rows[i].cells
            nombre = proyecto.get('nombre', proyecto.get('proyecto', 'Sin nombre'))
            score = proyecto.get('score_final', 0)

            # Determinar prioridad
            if score >= 70:
                prioridad = "ALTA"
            elif score >= 50:
                prioridad = "MEDIA"
            else:
                prioridad = "BAJA"

            row[0].text = str(i)
            row[1].text = nombre[:50] + "..." if len(nombre) > 50 else nombre
            row[2].text = f"{score:.1f}"
            row[3].text = prioridad

        doc.add_page_break()

    def _agregar_comparativo_word(self, doc: Document):
        """Agrega análisis comparativo por criterio."""
        doc.add_heading("4. Análisis Comparativo por Criterio", level=1)

        if not self.resultados:
            doc.add_paragraph("No hay datos para análisis comparativo.")
            return

        for criterio, peso in self.PESOS_CRITERIOS.items():
            doc.add_heading(f"4.{list(self.PESOS_CRITERIOS.keys()).index(criterio)+1} {criterio} (Peso: {peso*100:.0f}%)", level=2)

            # Extraer scores del criterio
            scores_criterio = []
            for r in self.resultados:
                nombre = r.get('nombre', r.get('proyecto', 'Sin nombre'))

                # Buscar score del criterio
                score_criterio = 0
                if 'detalle_criterios' in r:
                    for det in r['detalle_criterios']:
                        if det.get('criterio', '').upper().startswith(criterio.upper()):
                            score_criterio = det.get('score_base', 0)
                            break

                scores_criterio.append((nombre, score_criterio))

            # Ordenar por score
            scores_criterio.sort(key=lambda x: x[1], reverse=True)

            # Tabla simple
            table = doc.add_table(rows=min(5, len(scores_criterio)) + 1, cols=2)
            table.style = 'Table Grid'

            headers = table.rows[0].cells
            headers[0].text = "Proyecto"
            headers[1].text = f"Score {criterio}"
            for cell in headers:
                cell.paragraphs[0].runs[0].bold = True

            for i, (nombre, score) in enumerate(scores_criterio[:5], 1):
                row = table.rows[i].cells
                row[0].text = nombre[:40] + "..." if len(nombre) > 40 else nombre
                row[1].text = f"{score:.1f}"

            doc.add_paragraph()

        doc.add_page_break()

    def _agregar_detalle_proyectos_word(self, doc: Document):
        """Agrega detalle de cada proyecto."""
        doc.add_heading("5. Detalle por Proyecto", level=1)

        if not self.resultados:
            doc.add_paragraph("No hay proyectos para detallar.")
            return

        # Ordenar por score
        proyectos_ordenados = sorted(
            self.resultados,
            key=lambda x: x.get('score_final', 0),
            reverse=True
        )

        for i, proyecto in enumerate(proyectos_ordenados, 1):
            nombre = proyecto.get('nombre', proyecto.get('proyecto', 'Sin nombre'))
            score = proyecto.get('score_final', 0)

            # Título del proyecto
            doc.add_heading(f"5.{i} {nombre}", level=2)

            # Score final
            p = doc.add_paragraph()
            p.add_run("Score Final: ").bold = True
            run = p.add_run(f"{score:.1f}/100")
            if score >= 70:
                run.font.color.rgb = self.COLOR_SECUNDARIO
            elif score < 50:
                run.font.color.rgb = self.COLOR_ALERTA

            # Desglose por criterio
            if 'detalle_criterios' in proyecto:
                doc.add_paragraph()
                p = doc.add_paragraph()
                p.add_run("Desglose por Criterio:").bold = True

                for det in proyecto['detalle_criterios']:
                    criterio = det.get('criterio', 'N/A')
                    score_base = det.get('score_base', 0)
                    contribucion = det.get('contribucion_parcial', det.get('contribucion', 0))

                    p = doc.add_paragraph(style='List Bullet')
                    p.add_run(f"{criterio}: ")
                    p.add_run(f"Score {score_base:.1f} → Contribución {contribucion:.1f}")

            # Alertas si existen
            alertas = proyecto.get('alertas', [])
            if alertas:
                doc.add_paragraph()
                p = doc.add_paragraph()
                p.add_run("Alertas:").bold = True

                for alerta in alertas:
                    p = doc.add_paragraph(alerta, style='List Bullet')

            doc.add_paragraph()

        doc.add_page_break()

    def _agregar_recomendaciones_word(self, doc: Document):
        """Agrega sección de recomendaciones."""
        doc.add_heading("6. Recomendaciones", level=1)

        if not self.resultados:
            doc.add_paragraph("No hay datos para generar recomendaciones.")
            return

        # Analizar resultados
        scores = [r.get('score_final', 0) for r in self.resultados]
        score_promedio = sum(scores) / len(scores) if scores else 0

        alta = sum(1 for s in scores if s >= 70)
        media = sum(1 for s in scores if 50 <= s < 70)
        baja = sum(1 for s in scores if s < 50)

        doc.add_heading("6.1 Recomendaciones Generales", level=2)

        # Proyectos de alta prioridad
        if alta > 0:
            p = doc.add_paragraph()
            p.add_run(f"Aprobación Recomendada: ").bold = True
            p.add_run(
                f"Se identificaron {alta} proyectos de alta prioridad (score ≥70) "
                f"que se recomienda aprobar para ejecución."
            )

        # Proyectos de media prioridad
        if media > 0:
            p = doc.add_paragraph()
            p.add_run(f"Revisión Sugerida: ").bold = True
            p.add_run(
                f"{media} proyectos tienen prioridad media (score 50-69). "
                f"Se sugiere revisar con ajustes menores antes de aprobar."
            )

        # Proyectos de baja prioridad
        if baja > 0:
            p = doc.add_paragraph()
            p.add_run(f"Requieren Rediseño: ").bold = True
            p.add_run(
                f"{baja} proyectos tienen baja prioridad (score <50). "
                f"Se recomienda rediseñar o reconsiderar su viabilidad."
            )

        doc.add_paragraph()

        doc.add_heading("6.2 Proyectos Recomendados para Aprobación", level=2)

        # Listar top proyectos
        proyectos_top = sorted(
            self.resultados,
            key=lambda x: x.get('score_final', 0),
            reverse=True
        )

        proyectos_aprobacion = [p for p in proyectos_top if p.get('score_final', 0) >= 70]

        if proyectos_aprobacion:
            for proyecto in proyectos_aprobacion:
                nombre = proyecto.get('nombre', proyecto.get('proyecto', 'Sin nombre'))
                score = proyecto.get('score_final', 0)

                p = doc.add_paragraph(style='List Bullet')
                p.add_run(f"{nombre}").bold = True
                p.add_run(f" (Score: {score:.1f})")
        else:
            doc.add_paragraph(
                "No se identificaron proyectos con score suficiente para aprobación inmediata."
            )

        doc.add_paragraph()

        doc.add_heading("6.3 Consideraciones para el Comité", level=2)

        consideraciones = [
            "Revisar alertas específicas de cada proyecto antes de aprobar",
            "Considerar capacidad de ejecución simultánea de proyectos",
            "Evaluar sinergias entre proyectos del mismo territorio",
            "Verificar disponibilidad presupuestaria para proyectos aprobados"
        ]

        for c in consideraciones:
            doc.add_paragraph(c, style='List Bullet')

        # Pie de reporte
        doc.add_paragraph()
        doc.add_paragraph()

        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = p.add_run("— Fin del Reporte —")
        run.italic = True
        run.font.color.rgb = RGBColor(128, 128, 128)

        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = p.add_run(f"Generado: {self.fecha_generacion.strftime('%Y-%m-%d %H:%M')}")
        run.font.size = Pt(9)
        run.font.color.rgb = RGBColor(128, 128, 128)

    # ==================== MÉTODOS EXCEL ====================

    def _crear_hoja_resumen_excel(self, wb: Workbook):
        """Crea hoja de resumen ejecutivo."""
        ws = wb.create_sheet("Resumen Ejecutivo", 0)

        # Estilos
        header_fill = PatternFill(start_color="006699", end_color="006699", fill_type="solid")
        header_font = Font(color="FFFFFF", bold=True, size=12)
        title_font = Font(bold=True, size=14, color="006699")

        # Título
        ws['A1'] = "REPORTE DE EVALUACIÓN DE CARTERA"
        ws['A1'].font = Font(bold=True, size=16, color="006699")
        ws.merge_cells('A1:D1')

        ws['A2'] = f"Fecha: {self.fecha_generacion.strftime('%Y-%m-%d %H:%M')}"
        ws['A2'].font = Font(italic=True, size=10)

        # Métricas principales
        ws['A4'] = "MÉTRICAS PRINCIPALES"
        ws['A4'].font = title_font

        num_proyectos = len(self.resultados)
        if num_proyectos > 0:
            scores = [r.get('score_final', 0) for r in self.resultados]
            score_promedio = sum(scores) / len(scores)
            score_max = max(scores)
            score_min = min(scores)
            alta = sum(1 for s in scores if s >= 70)
            media = sum(1 for s in scores if 50 <= s < 70)
            baja = sum(1 for s in scores if s < 50)
        else:
            score_promedio = score_max = score_min = 0
            alta = media = baja = 0

        metricas = [
            ("Total Proyectos", num_proyectos),
            ("Score Promedio", f"{score_promedio:.1f}"),
            ("Score Máximo", f"{score_max:.1f}"),
            ("Score Mínimo", f"{score_min:.1f}"),
            ("Prioridad Alta (≥70)", alta),
            ("Prioridad Media (50-69)", media),
            ("Prioridad Baja (<50)", baja),
        ]

        for i, (metrica, valor) in enumerate(metricas, 5):
            ws[f'A{i}'] = metrica
            ws[f'B{i}'] = valor
            ws[f'A{i}'].font = Font(bold=True)

        # Ajustar anchos
        ws.column_dimensions['A'].width = 25
        ws.column_dimensions['B'].width = 15
        ws.column_dimensions['C'].width = 20
        ws.column_dimensions['D'].width = 15

    def _crear_hoja_ranking_excel(self, wb: Workbook):
        """Crea hoja con ranking completo."""
        ws = wb.create_sheet("Ranking")

        # Estilos
        header_fill = PatternFill(start_color="006699", end_color="006699", fill_type="solid")
        header_font = Font(color="FFFFFF", bold=True)
        alto_fill = PatternFill(start_color="C6EFCE", end_color="C6EFCE", fill_type="solid")
        medio_fill = PatternFill(start_color="FFEB9C", end_color="FFEB9C", fill_type="solid")
        bajo_fill = PatternFill(start_color="FFC7CE", end_color="FFC7CE", fill_type="solid")

        # Encabezados
        headers = ["Posición", "Proyecto", "Score Final", "Prioridad", "SROI", "Stakeholders", "Probabilidad", "Riesgos"]
        for col, header in enumerate(headers, 1):
            cell = ws.cell(row=1, column=col, value=header)
            cell.fill = header_fill
            cell.font = header_font
            cell.alignment = Alignment(horizontal='center')

        # Ordenar proyectos
        proyectos_ordenados = sorted(
            self.resultados,
            key=lambda x: x.get('score_final', 0),
            reverse=True
        )

        # Datos
        for i, proyecto in enumerate(proyectos_ordenados, 1):
            row = i + 1
            nombre = proyecto.get('nombre', proyecto.get('proyecto', 'Sin nombre'))
            score = proyecto.get('score_final', 0)

            # Determinar prioridad
            if score >= 70:
                prioridad = "ALTA"
                fill = alto_fill
            elif score >= 50:
                prioridad = "MEDIA"
                fill = medio_fill
            else:
                prioridad = "BAJA"
                fill = bajo_fill

            # Extraer scores por criterio
            scores_criterios = {'SROI': 0, 'Stakeholders': 0, 'Probabilidad': 0, 'Riesgos': 0}
            if 'detalle_criterios' in proyecto:
                for det in proyecto['detalle_criterios']:
                    criterio = det.get('criterio', '')
                    for key in scores_criterios:
                        if criterio.upper().startswith(key.upper()):
                            scores_criterios[key] = det.get('score_base', 0)
                            break

            ws.cell(row=row, column=1, value=i)
            ws.cell(row=row, column=2, value=nombre)
            ws.cell(row=row, column=3, value=round(score, 1))
            ws.cell(row=row, column=4, value=prioridad).fill = fill
            ws.cell(row=row, column=5, value=round(scores_criterios['SROI'], 1))
            ws.cell(row=row, column=6, value=round(scores_criterios['Stakeholders'], 1))
            ws.cell(row=row, column=7, value=round(scores_criterios['Probabilidad'], 1))
            ws.cell(row=row, column=8, value=round(scores_criterios['Riesgos'], 1))

        # Ajustar anchos
        ws.column_dimensions['A'].width = 10
        ws.column_dimensions['B'].width = 40
        ws.column_dimensions['C'].width = 12
        ws.column_dimensions['D'].width = 12
        ws.column_dimensions['E'].width = 10
        ws.column_dimensions['F'].width = 12
        ws.column_dimensions['G'].width = 12
        ws.column_dimensions['H'].width = 10

    def _crear_hoja_comparativo_excel(self, wb: Workbook):
        """Crea hoja con análisis comparativo por criterio."""
        ws = wb.create_sheet("Comparativo")

        # Estilos
        header_fill = PatternFill(start_color="006699", end_color="006699", fill_type="solid")
        header_font = Font(color="FFFFFF", bold=True)
        criterio_font = Font(bold=True, size=11, color="006699")

        row_actual = 1

        for criterio, peso in self.PESOS_CRITERIOS.items():
            # Título del criterio
            ws.cell(row=row_actual, column=1, value=f"{criterio} (Peso: {peso*100:.0f}%)")
            ws.cell(row=row_actual, column=1).font = criterio_font
            ws.merge_cells(start_row=row_actual, start_column=1, end_row=row_actual, end_column=3)

            row_actual += 1

            # Encabezados
            ws.cell(row=row_actual, column=1, value="Proyecto").fill = header_fill
            ws.cell(row=row_actual, column=1).font = header_font
            ws.cell(row=row_actual, column=2, value="Score").fill = header_fill
            ws.cell(row=row_actual, column=2).font = header_font
            ws.cell(row=row_actual, column=3, value="Contribución").fill = header_fill
            ws.cell(row=row_actual, column=3).font = header_font

            row_actual += 1

            # Extraer y ordenar scores del criterio
            scores_criterio = []
            for r in self.resultados:
                nombre = r.get('nombre', r.get('proyecto', 'Sin nombre'))
                score_base = 0
                contribucion = 0

                if 'detalle_criterios' in r:
                    for det in r['detalle_criterios']:
                        if det.get('criterio', '').upper().startswith(criterio.upper()):
                            score_base = det.get('score_base', 0)
                            contribucion = det.get('contribucion_parcial', det.get('contribucion', 0))
                            break

                scores_criterio.append((nombre, score_base, contribucion))

            scores_criterio.sort(key=lambda x: x[1], reverse=True)

            for nombre, score, contrib in scores_criterio:
                ws.cell(row=row_actual, column=1, value=nombre)
                ws.cell(row=row_actual, column=2, value=round(score, 1))
                ws.cell(row=row_actual, column=3, value=round(contrib, 2))
                row_actual += 1

            row_actual += 2  # Espacio entre criterios

        # Ajustar anchos
        ws.column_dimensions['A'].width = 40
        ws.column_dimensions['B'].width = 12
        ws.column_dimensions['C'].width = 15

    def _crear_hoja_detalle_excel(self, wb: Workbook):
        """Crea hoja con detalle de todos los proyectos."""
        ws = wb.create_sheet("Detalle Proyectos")

        # Estilos
        header_fill = PatternFill(start_color="006699", end_color="006699", fill_type="solid")
        header_font = Font(color="FFFFFF", bold=True)

        # Encabezados expandidos
        headers = [
            "Proyecto", "Score Final", "Prioridad",
            "SROI Score", "SROI Contrib",
            "Stakeholders Score", "Stakeholders Contrib",
            "Probabilidad Score", "Probabilidad Contrib",
            "Riesgos Score", "Riesgos Contrib",
            "Alertas"
        ]

        for col, header in enumerate(headers, 1):
            cell = ws.cell(row=1, column=col, value=header)
            cell.fill = header_fill
            cell.font = header_font
            cell.alignment = Alignment(horizontal='center', wrap_text=True)

        # Ordenar proyectos
        proyectos_ordenados = sorted(
            self.resultados,
            key=lambda x: x.get('score_final', 0),
            reverse=True
        )

        for i, proyecto in enumerate(proyectos_ordenados, 1):
            row = i + 1
            nombre = proyecto.get('nombre', proyecto.get('proyecto', 'Sin nombre'))
            score = proyecto.get('score_final', 0)

            # Prioridad
            if score >= 70:
                prioridad = "ALTA"
            elif score >= 50:
                prioridad = "MEDIA"
            else:
                prioridad = "BAJA"

            # Extraer scores y contribuciones
            criterios_data = {
                'SROI': {'score': 0, 'contrib': 0},
                'Stakeholders': {'score': 0, 'contrib': 0},
                'Probabilidad': {'score': 0, 'contrib': 0},
                'Riesgos': {'score': 0, 'contrib': 0}
            }

            if 'detalle_criterios' in proyecto:
                for det in proyecto['detalle_criterios']:
                    criterio = det.get('criterio', '')
                    for key in criterios_data:
                        if criterio.upper().startswith(key.upper()):
                            criterios_data[key]['score'] = det.get('score_base', 0)
                            criterios_data[key]['contrib'] = det.get('contribucion_parcial', det.get('contribucion', 0))
                            break

            # Alertas
            alertas = proyecto.get('alertas', [])
            alertas_texto = "; ".join(alertas) if alertas else ""

            # Escribir fila
            ws.cell(row=row, column=1, value=nombre)
            ws.cell(row=row, column=2, value=round(score, 1))
            ws.cell(row=row, column=3, value=prioridad)
            ws.cell(row=row, column=4, value=round(criterios_data['SROI']['score'], 1))
            ws.cell(row=row, column=5, value=round(criterios_data['SROI']['contrib'], 2))
            ws.cell(row=row, column=6, value=round(criterios_data['Stakeholders']['score'], 1))
            ws.cell(row=row, column=7, value=round(criterios_data['Stakeholders']['contrib'], 2))
            ws.cell(row=row, column=8, value=round(criterios_data['Probabilidad']['score'], 1))
            ws.cell(row=row, column=9, value=round(criterios_data['Probabilidad']['contrib'], 2))
            ws.cell(row=row, column=10, value=round(criterios_data['Riesgos']['score'], 1))
            ws.cell(row=row, column=11, value=round(criterios_data['Riesgos']['contrib'], 2))
            ws.cell(row=row, column=12, value=alertas_texto)

        # Ajustar anchos
        anchos = [40, 12, 10, 10, 12, 15, 15, 15, 15, 12, 12, 50]
        for i, ancho in enumerate(anchos, 1):
            ws.column_dimensions[ws.cell(row=1, column=i).column_letter].width = ancho

    def _crear_hoja_metodologia_excel(self, wb: Workbook):
        """Crea hoja explicando la metodología."""
        ws = wb.create_sheet("Metodología")

        # Estilos
        title_font = Font(bold=True, size=14, color="006699")
        header_fill = PatternFill(start_color="006699", end_color="006699", fill_type="solid")
        header_font = Font(color="FFFFFF", bold=True)

        # Título
        ws['A1'] = "METODOLOGÍA DE EVALUACIÓN - ARQUITECTURA C"
        ws['A1'].font = title_font
        ws.merge_cells('A1:D1')

        # Descripción
        ws['A3'] = "El Sistema de Priorización ENLAZA utiliza la Arquitectura C para evaluar proyectos sociales."
        ws.merge_cells('A3:D3')

        # Criterios y pesos
        ws['A5'] = "CRITERIOS DE EVALUACIÓN"
        ws['A5'].font = Font(bold=True, size=12)

        ws.cell(row=6, column=1, value="Criterio").fill = header_fill
        ws.cell(row=6, column=1).font = header_font
        ws.cell(row=6, column=2, value="Peso").fill = header_fill
        ws.cell(row=6, column=2).font = header_font
        ws.cell(row=6, column=3, value="Descripción").fill = header_fill
        ws.cell(row=6, column=3).font = header_font

        criterios = [
            ("SROI", "40%", "Retorno Social de la Inversión - Métrica dominante que determina el valor social generado"),
            ("Stakeholders", "25%", "Evalúa relacionamiento con comunidades y pertinencia operacional para ENLAZA"),
            ("Probabilidad", "20%", "Probabilidad de éxito basada en recursos, experiencia y condiciones"),
            ("Riesgos", "15%", "Evaluación de riesgos técnicos, sociales, financieros y regulatorios (scoring inverso)")
        ]

        for i, (criterio, peso, desc) in enumerate(criterios, 7):
            ws.cell(row=i, column=1, value=criterio)
            ws.cell(row=i, column=2, value=peso)
            ws.cell(row=i, column=3, value=desc)

        # Escalas
        ws['A13'] = "ESCALAS DE PRIORIZACIÓN"
        ws['A13'].font = Font(bold=True, size=12)

        ws.cell(row=14, column=1, value="Rango Score").fill = header_fill
        ws.cell(row=14, column=1).font = header_font
        ws.cell(row=14, column=2, value="Prioridad").fill = header_fill
        ws.cell(row=14, column=2).font = header_font
        ws.cell(row=14, column=3, value="Recomendación").fill = header_fill
        ws.cell(row=14, column=3).font = header_font

        escalas = [
            ("≥ 70", "ALTA", "Recomendado para aprobación inmediata"),
            ("50 - 69", "MEDIA", "Considerar con ajustes menores"),
            ("< 50", "BAJA", "Requiere rediseño o reconsideración")
        ]

        for i, (rango, prioridad, recom) in enumerate(escalas, 15):
            ws.cell(row=i, column=1, value=rango)
            ws.cell(row=i, column=2, value=prioridad)
            ws.cell(row=i, column=3, value=recom)

        # Consideraciones especiales
        ws['A20'] = "CONSIDERACIONES ESPECIALES"
        ws['A20'].font = Font(bold=True, size=12)

        consideraciones = [
            "• SROI < 1.0: Proyecto rechazado automáticamente (destruye valor social)",
            "• SROI > 7.0: Requiere verificación metodológica (valor atípico)",
            "• Riesgos críticos: Generan alertas especiales independiente del score",
            "• Municipios PDET: Bonificación por impacto en territorios priorizados"
        ]

        for i, cons in enumerate(consideraciones, 21):
            ws.cell(row=i, column=1, value=cons)
            ws.merge_cells(start_row=i, start_column=1, end_row=i, end_column=3)

        # Ajustar anchos
        ws.column_dimensions['A'].width = 20
        ws.column_dimensions['B'].width = 15
        ws.column_dimensions['C'].width = 60
        ws.column_dimensions['D'].width = 15

    # ==================== MÉTODOS POWERPOINT ====================

    def exportar_powerpoint(self) -> BytesIO:
        """
        Genera presentación PowerPoint ejecutiva para comité.

        Returns:
            BytesIO con el archivo PowerPoint
        """
        prs = Presentation()
        prs.slide_width = PptxInches(13.333)  # 16:9
        prs.slide_height = PptxInches(7.5)

        # Slide 1: Portada
        self._agregar_slide_portada_pptx(prs)

        # Slide 2: Resumen Ejecutivo
        self._agregar_slide_resumen_pptx(prs)

        # Slide 3: Ranking Top 5
        self._agregar_slide_ranking_pptx(prs)

        # Slide 4: Análisis por Criterio
        self._agregar_slide_criterios_pptx(prs)

        # Slide 5: Recomendaciones
        self._agregar_slide_recomendaciones_pptx(prs)

        # Guardar
        buffer = BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer

    def _agregar_slide_portada_pptx(self, prs: Presentation):
        """Slide de portada profesional."""
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)

        # Fondo azul corporativo (usando shape)
        shape = slide.shapes.add_shape(
            1,  # Rectangle
            PptxInches(0), PptxInches(0),
            PptxInches(13.333), PptxInches(2.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = PptxRGBColor(0, 102, 153)
        shape.line.fill.background()

        # Título ENLAZA
        titulo = slide.shapes.add_textbox(
            PptxInches(0.5), PptxInches(0.5),
            PptxInches(12), PptxInches(1)
        )
        tf = titulo.text_frame
        p = tf.paragraphs[0]
        p.text = "ENLAZA GEB"
        p.font.size = PptxPt(44)
        p.font.bold = True
        p.font.color.rgb = PptxRGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        # Subtítulo
        subtitulo = slide.shapes.add_textbox(
            PptxInches(0.5), PptxInches(1.5),
            PptxInches(12), PptxInches(0.7)
        )
        tf = subtitulo.text_frame
        p = tf.paragraphs[0]
        p.text = "Sistema de Priorización de Proyectos Sociales"
        p.font.size = PptxPt(20)
        p.font.color.rgb = PptxRGBColor(200, 220, 240)
        p.alignment = PP_ALIGN.CENTER

        # Título principal
        titulo_principal = slide.shapes.add_textbox(
            PptxInches(0.5), PptxInches(3.2),
            PptxInches(12), PptxInches(1.2)
        )
        tf = titulo_principal.text_frame
        p = tf.paragraphs[0]
        p.text = "EVALUACIÓN DE CARTERA"
        p.font.size = PptxPt(48)
        p.font.bold = True
        p.font.color.rgb = PptxRGBColor(0, 102, 153)
        p.alignment = PP_ALIGN.CENTER

        # Información
        num_proyectos = len(self.resultados)
        info = slide.shapes.add_textbox(
            PptxInches(0.5), PptxInches(4.5),
            PptxInches(12), PptxInches(0.5)
        )
        tf = info.text_frame
        p = tf.paragraphs[0]
        p.text = f"{num_proyectos} proyectos evaluados | Metodología Arquitectura C"
        p.font.size = PptxPt(18)
        p.font.color.rgb = PptxRGBColor(100, 100, 100)
        p.alignment = PP_ALIGN.CENTER

        # Fecha
        fecha = slide.shapes.add_textbox(
            PptxInches(0.5), PptxInches(6.5),
            PptxInches(12), PptxInches(0.5)
        )
        tf = fecha.text_frame
        p = tf.paragraphs[0]
        p.text = self.fecha_generacion.strftime("%d de %B de %Y")
        p.font.size = PptxPt(14)
        p.font.color.rgb = PptxRGBColor(128, 128, 128)
        p.alignment = PP_ALIGN.CENTER

    def _agregar_slide_resumen_pptx(self, prs: Presentation):
        """Slide de resumen ejecutivo con métricas."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # Título
        titulo = slide.shapes.add_textbox(
            PptxInches(0.5), PptxInches(0.3),
            PptxInches(12), PptxInches(0.8)
        )
        tf = titulo.text_frame
        p = tf.paragraphs[0]
        p.text = "RESUMEN EJECUTIVO"
        p.font.size = PptxPt(36)
        p.font.bold = True
        p.font.color.rgb = PptxRGBColor(0, 102, 153)

        # Calcular métricas
        num_proyectos = len(self.resultados)
        if num_proyectos > 0:
            scores = [r.get('score_final', 0) for r in self.resultados]
            score_promedio = sum(scores) / len(scores)
            alta = sum(1 for s in scores if s >= 70)
            media = sum(1 for s in scores if 50 <= s < 70)
            baja = sum(1 for s in scores if s < 50)
        else:
            score_promedio = 0
            alta = media = baja = 0

        # Métricas en cajas
        metricas = [
            ("Total Proyectos", str(num_proyectos), PptxRGBColor(0, 102, 153)),
            ("Score Promedio", f"{score_promedio:.1f}", PptxRGBColor(0, 153, 102)),
            ("Alta Prioridad", str(alta), PptxRGBColor(46, 125, 50)),
            ("Media Prioridad", str(media), PptxRGBColor(255, 152, 0)),
            ("Baja Prioridad", str(baja), PptxRGBColor(211, 47, 47))
        ]

        for i, (label, valor, color) in enumerate(metricas):
            x = PptxInches(0.5 + i * 2.5)
            y = PptxInches(1.5)

            # Caja
            box = slide.shapes.add_shape(
                1, x, y, PptxInches(2.3), PptxInches(1.8)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = color
            box.line.fill.background()

            # Valor
            valor_box = slide.shapes.add_textbox(x, PptxInches(1.7), PptxInches(2.3), PptxInches(0.8))
            tf = valor_box.text_frame
            p = tf.paragraphs[0]
            p.text = valor
            p.font.size = PptxPt(36)
            p.font.bold = True
            p.font.color.rgb = PptxRGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER

            # Label
            label_box = slide.shapes.add_textbox(x, PptxInches(2.5), PptxInches(2.3), PptxInches(0.6))
            tf = label_box.text_frame
            p = tf.paragraphs[0]
            p.text = label
            p.font.size = PptxPt(12)
            p.font.color.rgb = PptxRGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER

        # Distribución textual
        dist_text = slide.shapes.add_textbox(
            PptxInches(0.5), PptxInches(4),
            PptxInches(12), PptxInches(2)
        )
        tf = dist_text.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = "Distribución de Proyectos por Nivel de Prioridad:"
        p.font.size = PptxPt(16)
        p.font.bold = True
        p.space_after = PptxPt(12)

        if num_proyectos > 0:
            textos = [
                f"• Alta Prioridad (≥70): {alta} proyectos ({alta/num_proyectos*100:.0f}%) - Recomendados para aprobación",
                f"• Media Prioridad (50-69): {media} proyectos ({media/num_proyectos*100:.0f}%) - Requieren revisión",
                f"• Baja Prioridad (<50): {baja} proyectos ({baja/num_proyectos*100:.0f}%) - Requieren rediseño"
            ]
            for texto in textos:
                p = tf.add_paragraph()
                p.text = texto
                p.font.size = PptxPt(14)
                p.space_after = PptxPt(8)

    def _agregar_slide_ranking_pptx(self, prs: Presentation):
        """Slide con ranking Top 5."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # Título
        titulo = slide.shapes.add_textbox(
            PptxInches(0.5), PptxInches(0.3),
            PptxInches(12), PptxInches(0.8)
        )
        tf = titulo.text_frame
        p = tf.paragraphs[0]
        p.text = "TOP 5 PROYECTOS RECOMENDADOS"
        p.font.size = PptxPt(36)
        p.font.bold = True
        p.font.color.rgb = PptxRGBColor(0, 102, 153)

        # Ordenar proyectos
        proyectos_ordenados = sorted(
            self.resultados,
            key=lambda x: x.get('score_final', 0),
            reverse=True
        )[:5]

        medallas = ['1', '2', '3', '4', '5']
        colores_medalla = [
            PptxRGBColor(255, 215, 0),   # Oro
            PptxRGBColor(192, 192, 192), # Plata
            PptxRGBColor(205, 127, 50),  # Bronce
            PptxRGBColor(0, 102, 153),   # Azul
            PptxRGBColor(0, 102, 153)    # Azul
        ]

        for i, proyecto in enumerate(proyectos_ordenados):
            nombre = proyecto.get('nombre', proyecto.get('proyecto', 'Sin nombre'))
            score = proyecto.get('score_final', 0)

            y = PptxInches(1.3 + i * 1.1)

            # Número/medalla
            num_box = slide.shapes.add_shape(
                9,  # Oval
                PptxInches(0.5), y,
                PptxInches(0.7), PptxInches(0.7)
            )
            num_box.fill.solid()
            num_box.fill.fore_color.rgb = colores_medalla[i]
            num_box.line.fill.background()

            num_text = slide.shapes.add_textbox(
                PptxInches(0.5), PptxInches(y.inches + 0.15),
                PptxInches(0.7), PptxInches(0.5)
            )
            tf = num_text.text_frame
            p = tf.paragraphs[0]
            p.text = medallas[i]
            p.font.size = PptxPt(24)
            p.font.bold = True
            p.font.color.rgb = PptxRGBColor(255, 255, 255) if i < 3 else PptxRGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER

            # Nombre del proyecto
            nombre_box = slide.shapes.add_textbox(
                PptxInches(1.4), y,
                PptxInches(8), PptxInches(0.7)
            )
            tf = nombre_box.text_frame
            p = tf.paragraphs[0]
            p.text = nombre[:60] + "..." if len(nombre) > 60 else nombre
            p.font.size = PptxPt(18)
            p.font.bold = True
            p.font.color.rgb = PptxRGBColor(50, 50, 50)

            # Score
            score_box = slide.shapes.add_textbox(
                PptxInches(10), y,
                PptxInches(2), PptxInches(0.7)
            )
            tf = score_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"{score:.1f}/100"
            p.font.size = PptxPt(20)
            p.font.bold = True
            if score >= 70:
                p.font.color.rgb = PptxRGBColor(46, 125, 50)
            elif score >= 50:
                p.font.color.rgb = PptxRGBColor(255, 152, 0)
            else:
                p.font.color.rgb = PptxRGBColor(211, 47, 47)
            p.alignment = PP_ALIGN.RIGHT

    def _agregar_slide_criterios_pptx(self, prs: Presentation):
        """Slide con análisis por criterio."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # Título
        titulo = slide.shapes.add_textbox(
            PptxInches(0.5), PptxInches(0.3),
            PptxInches(12), PptxInches(0.8)
        )
        tf = titulo.text_frame
        p = tf.paragraphs[0]
        p.text = "METODOLOGÍA DE EVALUACIÓN"
        p.font.size = PptxPt(36)
        p.font.bold = True
        p.font.color.rgb = PptxRGBColor(0, 102, 153)

        # Criterios
        criterios = [
            ("SROI", "40%", "Retorno Social de la Inversión", PptxRGBColor(0, 102, 153)),
            ("Stakeholders", "25%", "Relacionamiento y Pertinencia", PptxRGBColor(0, 153, 102)),
            ("Probabilidad", "20%", "Probabilidad de Éxito", PptxRGBColor(255, 152, 0)),
            ("Riesgos", "15%", "Evaluación de Riesgos (Inverso)", PptxRGBColor(211, 47, 47))
        ]

        for i, (nombre, peso, desc, color) in enumerate(criterios):
            y = PptxInches(1.5 + i * 1.4)

            # Barra de peso
            ancho_barra = float(peso.replace('%', '')) / 100 * 10
            barra = slide.shapes.add_shape(
                1, PptxInches(0.5), y,
                PptxInches(ancho_barra), PptxInches(0.8)
            )
            barra.fill.solid()
            barra.fill.fore_color.rgb = color
            barra.line.fill.background()

            # Nombre y peso
            texto = slide.shapes.add_textbox(
                PptxInches(0.7), PptxInches(y.inches + 0.15),
                PptxInches(4), PptxInches(0.5)
            )
            tf = texto.text_frame
            p = tf.paragraphs[0]
            p.text = f"{nombre} ({peso})"
            p.font.size = PptxPt(18)
            p.font.bold = True
            p.font.color.rgb = PptxRGBColor(255, 255, 255)

            # Descripción
            desc_box = slide.shapes.add_textbox(
                PptxInches(0.5), PptxInches(y.inches + 0.85),
                PptxInches(10), PptxInches(0.4)
            )
            tf = desc_box.text_frame
            p = tf.paragraphs[0]
            p.text = desc
            p.font.size = PptxPt(12)
            p.font.color.rgb = PptxRGBColor(100, 100, 100)

    def _agregar_slide_recomendaciones_pptx(self, prs: Presentation):
        """Slide de recomendaciones al comité."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # Título
        titulo = slide.shapes.add_textbox(
            PptxInches(0.5), PptxInches(0.3),
            PptxInches(12), PptxInches(0.8)
        )
        tf = titulo.text_frame
        p = tf.paragraphs[0]
        p.text = "RECOMENDACIONES AL COMITÉ"
        p.font.size = PptxPt(36)
        p.font.bold = True
        p.font.color.rgb = PptxRGBColor(0, 102, 153)

        # Calcular métricas
        if self.resultados:
            scores = [r.get('score_final', 0) for r in self.resultados]
            alta = sum(1 for s in scores if s >= 70)
            media = sum(1 for s in scores if 50 <= s < 70)
            baja = sum(1 for s in scores if s < 50)
        else:
            alta = media = baja = 0

        # Recomendaciones
        recomendaciones = [
            (f"APROBAR: {alta} proyecto(s)", "Proyectos con alta prioridad listos para ejecución", PptxRGBColor(46, 125, 50)),
            (f"REVISAR: {media} proyecto(s)", "Requieren análisis adicional o ajustes menores", PptxRGBColor(255, 152, 0)),
            (f"REPLANTEAR: {baja} proyecto(s)", "Necesitan rediseño significativo antes de aprobación", PptxRGBColor(211, 47, 47))
        ]

        for i, (titulo_rec, desc, color) in enumerate(recomendaciones):
            y = PptxInches(1.5 + i * 1.6)

            # Icono/indicador
            indicador = slide.shapes.add_shape(
                1, PptxInches(0.5), y,
                PptxInches(0.3), PptxInches(1.2)
            )
            indicador.fill.solid()
            indicador.fill.fore_color.rgb = color
            indicador.line.fill.background()

            # Título recomendación
            titulo_box = slide.shapes.add_textbox(
                PptxInches(1), y,
                PptxInches(10), PptxInches(0.6)
            )
            tf = titulo_box.text_frame
            p = tf.paragraphs[0]
            p.text = titulo_rec
            p.font.size = PptxPt(24)
            p.font.bold = True
            p.font.color.rgb = color

            # Descripción
            desc_box = slide.shapes.add_textbox(
                PptxInches(1), PptxInches(y.inches + 0.6),
                PptxInches(10), PptxInches(0.5)
            )
            tf = desc_box.text_frame
            p = tf.paragraphs[0]
            p.text = desc
            p.font.size = PptxPt(14)
            p.font.color.rgb = PptxRGBColor(100, 100, 100)

        # Nota final
        nota = slide.shapes.add_textbox(
            PptxInches(0.5), PptxInches(6.3),
            PptxInches(12), PptxInches(0.5)
        )
        tf = nota.text_frame
        p = tf.paragraphs[0]
        p.text = f"Reporte generado: {self.fecha_generacion.strftime('%Y-%m-%d %H:%M')} | Sistema de Priorización ENLAZA"
        p.font.size = PptxPt(10)
        p.font.color.rgb = PptxRGBColor(150, 150, 150)
        p.alignment = PP_ALIGN.CENTER
